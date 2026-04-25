import streamlit as st
import json
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import google.generativeai as genai
import os
import io
import base64
import re

st.set_page_config(
    page_title="Deckorate AI",
    page_icon="🎨",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ─── Custom CSS ───────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600&display=swap');

/* Sidebar */
[data-testid="stSidebar"] {
    background: #0f1117;
    border-right: 1px solid #1e2130;
}
[data-testid="stSidebar"] * { color: #e8eaf0 !important; }
[data-testid="stSidebar"] .step-label {
    font-size: 11px;
    font-weight: 600;
    letter-spacing: 0.08em;
    text-transform: uppercase;
    color: #6b7280 !important;
    margin: 1.2rem 0 0.3rem;
}
[data-testid="stSidebar"] .step-active {
    background: #1e2d4a;
    border-left: 3px solid #4f8ef7;
    border-radius: 4px;
    padding: 6px 10px;
    margin: 2px 0;
    font-weight: 500;
}
[data-testid="stSidebar"] .step-done {
    color: #34d399 !important;
}
[data-testid="stSidebar"] .step-pending {
    color: #4b5563 !important;
}

/* Main area */
.main-header {
    font-size: 1.8rem;
    font-weight: 700;
    margin-bottom: 0.2rem;
}
.sub-header {
    color: #6b7280;
    font-size: 0.95rem;
    margin-bottom: 1.5rem;
}
.slide-card {
    border: 1px solid #e5e7eb;
    border-radius: 10px;
    padding: 12px;
    cursor: pointer;
    transition: all 0.15s;
    background: white;
}
.slide-card:hover { border-color: #4f8ef7; box-shadow: 0 0 0 2px #dbeafe; }
.slide-card.selected { border-color: #4f8ef7; box-shadow: 0 0 0 3px #bfdbfe; }
.template-pill {
    display: inline-block;
    background: #eff6ff;
    color: #1d4ed8;
    border-radius: 20px;
    padding: 2px 10px;
    font-size: 12px;
    font-weight: 500;
    margin: 2px;
    cursor: pointer;
}
.template-pill.active { background: #1d4ed8; color: white; }
.palette-swatch {
    display: inline-block;
    width: 28px;
    height: 28px;
    border-radius: 6px;
    margin: 2px;
    cursor: pointer;
    border: 2px solid transparent;
}
.palette-swatch.selected { border-color: #1d4ed8; }
.ai-badge {
    background: linear-gradient(135deg, #667eea, #764ba2);
    color: white;
    border-radius: 20px;
    padding: 2px 10px;
    font-size: 11px;
    font-weight: 600;
}
.change-tag {
    background: #fef3c7;
    color: #92400e;
    border-radius: 4px;
    padding: 1px 6px;
    font-size: 11px;
}
</style>
""", unsafe_allow_html=True)

# ─── Constants ────────────────────────────────────────────────────────────────
FONT_PAIRS = {
    "Georgia + Calibri (McKinsey classic)": ("Georgia", "Calibri"),
    "Times New Roman + Calibri (Formal/Banking)": ("Times New Roman", "Calibri"),
    "Trebuchet MS + Calibri Light (Modern BCG)": ("Trebuchet MS", "Calibri Light"),
    "Arial Black + Arial (Bold & Direct)": ("Arial Black", "Arial"),
    "Cambria + Calibri (Refined Academic)": ("Cambria", "Calibri"),
    "Palatino + Garamond (Premium Editorial)": ("Palatino Linotype", "Garamond"),
    "Calibri + Calibri Light (Minimal)": ("Calibri", "Calibri Light"),
}

PALETTES = {
    "McKinsey Navy": {"primary": "1a2e5a", "accent": "0072ce", "bg": "ffffff", "text": "1a1a2e"},
    "BCG Green":     {"primary": "005f4e", "accent": "00a878", "bg": "ffffff", "text": "1a2e1a"},
    "Bain Red":      {"primary": "cc0000", "accent": "ff4444", "bg": "ffffff", "text": "1a0000"},
    "Minimal Gray":  {"primary": "2d2d2d", "accent": "6b7280", "bg": "fafafa", "text": "111111"},
    "Deep Blue":     {"primary": "0f172a", "accent": "3b82f6", "bg": "ffffff", "text": "0f172a"},
    "Slate Modern":  {"primary": "334155", "accent": "8b5cf6", "bg": "f8fafc", "text": "1e293b"},
}

LAYOUT_TEMPLATES = {
    "title":   ["Full bleed title", "Centered minimal", "Left-aligned executive", "Dark hero"],
    "bullets": ["Clean bullets", "Two column", "Icon grid", "Big quote callout"],
    "data":    ["Chart focus", "Stat callouts", "Table + insight", "Before / after"],
    "agenda":  ["Numbered list", "Card grid", "Timeline", "Section dividers"],
    "closing": ["Bold CTA", "Next steps", "Thank you minimal", "Contact card"],
}

# ─── Session state ─────────────────────────────────────────────────────────────
def init_state():
    defaults = {
        "step": 1,
        "pptx_bytes": None,
        "slides_data": [],
        "selected_slide": 0,
        "font_pair": list(FONT_PAIRS.keys())[0],
        "palette_name": "McKinsey Navy",
        "custom_palette": None,
        "beautified": False,
        "beautified_pptx": None,
        "slide_templates": {},
        "slide_edits": {},
        "ai_suggestions": {},
    }
    for k, v in defaults.items():
        if k not in st.session_state:
            st.session_state[k] = v

init_state()

# ─── Helpers ──────────────────────────────────────────────────────────────────
def extract_slides(pptx_bytes):
    prs = Presentation(io.BytesIO(pptx_bytes))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        title = texts[0] if texts else f"Slide {i+1}"
        slides.append({
            "index": i,
            "title": title,
            "texts": texts,
            "shape_count": len(slide.shapes),
        })
    return slides

def detect_slide_type(slide):
    texts = slide["texts"]
    all_text = " ".join(texts).lower()
    if slide["index"] == 0:
        return "title"
    if any(w in all_text for w in ["agenda", "contents", "overview", "outline"]):
        return "agenda"
    if any(w in all_text for w in ["%", "growth", "revenue", "cagr", "$", "million", "billion"]):
        return "data"
    if slide["index"] == len(st.session_state.slides_data) - 1:
        return "closing"
    return "bullets"

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def get_ai_suggestion(slide):
    """Ask Gemini to suggest layout + typo fixes for a slide."""
    api_key = st.session_state.get("gemini_api_key") or os.environ.get("GEMINI_API_KEY", "")
    if not api_key:
        return {
            "suggested_layout": "Clean bullets",
            "layout_reason": "Add your Gemini API key in the sidebar to get AI suggestions.",
            "typo_fixes": [],
            "design_tips": ["Use consistent font sizes", "Left-align body text"]
        }
    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel("gemini-1.5-flash")
        texts_str = "\n".join(f"- {t}" for t in slide["texts"][:10])
        prompt = f"""You are a consulting slide design expert. Analyze this slide and respond ONLY with JSON, no markdown.

Slide title: {slide['title']}
Slide content:
{texts_str}

Respond with this exact JSON structure:
{{
  "suggested_layout": "one of: Clean bullets | Two column | Icon grid | Big quote callout | Chart focus | Stat callouts | Full bleed title | Centered minimal | Numbered list | Card grid | Bold CTA | Next steps",
  "layout_reason": "one sentence why",
  "typo_fixes": [{{"original": "...", "fixed": "..."}}],
  "design_tips": ["tip 1", "tip 2"]
}}"""
        response = model.generate_content(prompt)
        raw = response.text.strip()
        raw = re.sub(r"```json|```", "", raw).strip()
        return json.loads(raw)
    except Exception:
        return {
            "suggested_layout": "Clean bullets",
            "layout_reason": "Default layout for text-heavy slides.",
            "typo_fixes": [],
            "design_tips": ["Use consistent font sizes", "Left-align body text"]
        }

def apply_beautification(pptx_bytes, palette, font_pair_key, slide_templates, slide_edits):
    """Apply palette, fonts, and per-slide edits to the pptx."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    pal = palette
    heading_font, body_font = FONT_PAIRS[font_pair_key]
    primary_rgb   = hex_to_rgb(pal["primary"])
    accent_rgb    = hex_to_rgb(pal["accent"])
    text_rgb      = hex_to_rgb(pal["text"])

    for i, slide in enumerate(prs.slides):
        edits = slide_edits.get(i, {})
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                for run in para.runs:
                    # Apply font
                    if shape.name.lower().startswith("title") or para_idx == 0:
                        run.font.name = heading_font
                        run.font.size = Pt(edits.get("heading_size", 32))
                        run.font.color.rgb = primary_rgb
                        run.font.bold = True
                    else:
                        run.font.name = body_font
                        run.font.size = Pt(edits.get("body_size", 16))
                        run.font.color.rgb = text_rgb

                    # Apply text edits
                    if "text_edits" in edits:
                        for orig, fixed in edits["text_edits"].items():
                            if orig in run.text:
                                run.text = run.text.replace(orig, fixed)

                # Alignment
                para.alignment = PP_ALIGN.LEFT

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

# ─── Sidebar ──────────────────────────────────────────────────────────────────
def render_sidebar():
    with st.sidebar:
        st.markdown("## 🎨 Deckorate AI")
        st.markdown("---")

        steps = [
            (1, "Upload deck"),
            (2, "Pick palette"),
            (3, "Choose fonts"),
            (4, "Cleanup options"),
            (5, "Preview & edit"),
            (6, "Beautify"),
            (7, "Download"),
        ]

        current = st.session_state.step
        for num, label in steps:
            if num < current:
                icon = "✅"
                style = "step-done"
            elif num == current:
                icon = "▶"
                style = "step-active"
            else:
                icon = "○"
                style = "step-pending"
            st.markdown(f'<div class="{style}">{icon} &nbsp; <b>{num}.</b> {label}</div>', unsafe_allow_html=True)

        st.markdown("---")
        st.markdown("**Gemini API key**")
        api_key_input = st.text_input(
            "Gemini API key",
            type="password",
            placeholder="AIza...",
            label_visibility="collapsed",
            value=st.session_state.get("gemini_api_key", "")
        )
        if api_key_input:
            st.session_state["gemini_api_key"] = api_key_input
            st.markdown('<span style="color:#34d399;font-size:12px;">✓ Key saved</span>', unsafe_allow_html=True)
        else:
            st.markdown('<span style="color:#6b7280;font-size:11px;">Get a free key at ai.google.dev</span>', unsafe_allow_html=True)
        st.markdown("---")
        if st.session_state.slides_data:
            st.markdown(f"**{len(st.session_state.slides_data)} slides** loaded")
        if st.session_state.beautified:
            st.markdown("✨ **Beautified!**")

# ─── Step renderers ───────────────────────────────────────────────────────────
def step_upload():
    st.markdown('<div class="main-header">Upload your deck</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">Accepts .pptx files. Google Slides: File → Download → PowerPoint first.</div>', unsafe_allow_html=True)

    uploaded = st.file_uploader("Drop your .pptx here", type=["pptx"], label_visibility="collapsed")
    if uploaded:
        st.session_state.pptx_bytes = uploaded.read()
        st.session_state.slides_data = extract_slides(st.session_state.pptx_bytes)
        st.success(f"✅ Loaded **{len(st.session_state.slides_data)} slides** from *{uploaded.name}*")

        # Preview slide titles
        st.markdown("**Slides detected:**")
        for s in st.session_state.slides_data:
            st.markdown(f"&nbsp;&nbsp;`{s['index']+1}` &nbsp; {s['title']}")

        if st.button("Continue →", type="primary"):
            st.session_state.step = 2
            st.rerun()

def step_palette():
    st.markdown('<div class="main-header">Pick a colour palette</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">Choose a preset consulting palette or build your own.</div>', unsafe_allow_html=True)

    col1, col2 = st.columns([1, 1])

    with col1:
        st.markdown("#### Preset palettes")
        for name, pal in PALETTES.items():
            selected = st.session_state.palette_name == name
            border = "3px solid #1d4ed8" if selected else "1px solid #e5e7eb"
            if st.button(f"{'✓ ' if selected else ''}{name}", key=f"pal_{name}",
                         use_container_width=True):
                st.session_state.palette_name = name
                st.session_state.custom_palette = None
                st.rerun()

    with col2:
        st.markdown("#### Palette preview")
        pal = st.session_state.custom_palette or PALETTES[st.session_state.palette_name]
        st.markdown(f"""
        <div style="border-radius:10px;overflow:hidden;border:1px solid #e5e7eb;">
          <div style="background:#{pal['primary']};padding:16px 20px;color:white;font-weight:700;font-size:1.1rem;">Slide Title Here</div>
          <div style="background:#{pal['bg']};padding:16px 20px;">
            <div style="color:#{pal['accent']};font-weight:600;margin-bottom:6px;">Key insight or section header</div>
            <div style="color:#{pal['text']};font-size:14px;">• Body text looks like this on your slides<br>• Second bullet point with supporting detail<br>• Third point to round out the argument</div>
          </div>
        </div>
        """, unsafe_allow_html=True)

        st.markdown("#### Or build your own")
        c1, c2, c3 = st.columns(3)
        with c1:
            p = st.color_picker("Primary", f"#{PALETTES[st.session_state.palette_name]['primary']}")
        with c2:
            a = st.color_picker("Accent", f"#{PALETTES[st.session_state.palette_name]['accent']}")
        with c3:
            b = st.color_picker("Background", f"#{PALETTES[st.session_state.palette_name]['bg']}")
        if st.button("Use custom palette"):
            st.session_state.custom_palette = {
                "primary": p.lstrip("#"), "accent": a.lstrip("#"),
                "bg": b.lstrip("#"), "text": "111111"
            }
            st.rerun()

    col_back, col_next = st.columns([1,5])
    with col_back:
        if st.button("← Back"): st.session_state.step = 1; st.rerun()
    with col_next:
        if st.button("Continue →", type="primary"): st.session_state.step = 3; st.rerun()

def step_fonts():
    st.markdown('<div class="main-header">Choose your fonts</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">All fonts are system-safe and will render correctly in PowerPoint on any machine.</div>', unsafe_allow_html=True)

    chosen = st.selectbox("Font pairing", list(FONT_PAIRS.keys()),
                          index=list(FONT_PAIRS.keys()).index(st.session_state.font_pair))
    st.session_state.font_pair = chosen
    heading_font, body_font = FONT_PAIRS[chosen]

    pal = st.session_state.custom_palette or PALETTES[st.session_state.palette_name]

    st.markdown("#### Live preview")
    st.markdown(f"""
    <div style="border:1px solid #e5e7eb;border-radius:10px;padding:24px 28px;background:#{pal['bg']};">
      <div style="font-family:'{heading_font}',serif;font-size:1.7rem;font-weight:700;color:#{pal['primary']};margin-bottom:8px;">Strategic Market Entry Analysis</div>
      <div style="font-family:'{heading_font}',serif;font-size:1rem;font-weight:600;color:#{pal['accent']};margin-bottom:12px;">Executive Summary</div>
      <div style="font-family:'{body_font}',sans-serif;font-size:0.9rem;color:#{pal['text']};line-height:1.7;">
        • The addressable market represents a $4.2B opportunity growing at 18% CAGR<br>
        • Three strategic vectors identified for sustainable competitive advantage<br>
        • Recommended investment of $12M over 24 months with 3.1x projected ROI
      </div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown(f"**Heading:** {heading_font} &nbsp;|&nbsp; **Body:** {body_font}")

    col_back, col_next = st.columns([1,5])
    with col_back:
        if st.button("← Back"): st.session_state.step = 2; st.rerun()
    with col_next:
        if st.button("Continue →", type="primary"): st.session_state.step = 4; st.rerun()

def step_cleanup():
    st.markdown('<div class="main-header">Cleanup options</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">Choose what the AI should fix across all slides.</div>', unsafe_allow_html=True)

    c1, c2 = st.columns(2)
    with c1:
        st.markdown("#### Text")
        fix_typos    = st.toggle("Fix typos & grammar", value=True)
        fix_case     = st.toggle("Fix inconsistent casing", value=True)
        fix_bullets  = st.toggle("Clean up bullet formatting", value=True)
    with c2:
        st.markdown("#### Layout")
        fix_align    = st.toggle("Standardise text alignment", value=True)
        fix_fonts    = st.toggle("Enforce consistent font sizes", value=True)
        fix_spacing  = st.toggle("Fix uneven spacing", value=True)

    st.markdown("#### Font sizes")
    col_h, col_b = st.columns(2)
    with col_h:
        heading_size = st.slider("Heading size (pt)", 24, 48, 32)
    with col_b:
        body_size = st.slider("Body text size (pt)", 12, 24, 16)

    # Store globally
    st.session_state["cleanup_opts"] = {
        "fix_typos": fix_typos, "fix_case": fix_case,
        "fix_bullets": fix_bullets, "fix_align": fix_align,
        "fix_fonts": fix_fonts, "fix_spacing": fix_spacing,
        "heading_size": heading_size, "body_size": body_size,
    }

    col_back, col_next = st.columns([1,5])
    with col_back:
        if st.button("← Back"): st.session_state.step = 3; st.rerun()
    with col_next:
        if st.button("Continue →", type="primary"): st.session_state.step = 5; st.rerun()

def step_preview():
    st.markdown('<div class="main-header">Preview & edit slides</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">Click any slide to edit text, swap layout, or change per-slide fonts and colours.</div>', unsafe_allow_html=True)

    slides = st.session_state.slides_data
    if not slides:
        st.warning("No slides loaded. Go back to step 1.")
        return

    # Thumbnail strip
    st.markdown("#### All slides")
    cols = st.columns(min(len(slides), 5))
    for i, slide in enumerate(slides):
        col = cols[i % 5]
        with col:
            slide_type = detect_slide_type(slide)
            is_selected = st.session_state.selected_slide == i
            border = "2px solid #4f8ef7" if is_selected else "1px solid #e5e7eb"
            bg = "#eff6ff" if is_selected else "white"
            if st.button(f"**{i+1}**\n{slide['title'][:20]}{'…' if len(slide['title'])>20 else ''}",
                         key=f"thumb_{i}", use_container_width=True):
                st.session_state.selected_slide = i
                # Get AI suggestion if not already fetched
                if i not in st.session_state.ai_suggestions:
                    with st.spinner("AI analysing slide…"):
                        st.session_state.ai_suggestions[i] = get_ai_suggestion(slide)
                st.rerun()

    st.markdown("---")

    # Detailed editor for selected slide
    idx = st.session_state.selected_slide
    slide = slides[idx]
    slide_type = detect_slide_type(slide)
    ai = st.session_state.ai_suggestions.get(idx)

    st.markdown(f"### Editing: Slide {idx+1} — {slide['title']}")

    left, right = st.columns([1.2, 1])

    with left:
        # AI suggestion banner
        if ai:
            st.markdown(f"""
            <div style="background:#f0f9ff;border:1px solid #bae6fd;border-radius:8px;padding:12px 16px;margin-bottom:12px;">
              <span class="ai-badge">AI suggestion</span>
              &nbsp;<b>{ai['suggested_layout']}</b> — {ai['layout_reason']}
            </div>
            """, unsafe_allow_html=True)

        # Layout template picker
        st.markdown("#### Layout template")
        templates = LAYOUT_TEMPLATES.get(slide_type, LAYOUT_TEMPLATES["bullets"])
        current_template = st.session_state.slide_templates.get(idx,
            ai["suggested_layout"] if ai else templates[0])

        template_cols = st.columns(2)
        for ti, tmpl in enumerate(templates):
            with template_cols[ti % 2]:
                is_active = current_template == tmpl
                label = f"✓ {tmpl}" if is_active else tmpl
                if st.button(label, key=f"tmpl_{idx}_{ti}", use_container_width=True,
                             type="primary" if is_active else "secondary"):
                    st.session_state.slide_templates[idx] = tmpl
                    st.rerun()

        # Per-slide font override
        with st.expander("Override fonts for this slide"):
            slide_heading = st.selectbox("Heading font", list(FONT_PAIRS.keys()),
                                          key=f"sfont_{idx}")
            st.session_state.setdefault("slide_font_overrides", {})[idx] = slide_heading

        # Per-slide color override
        with st.expander("Override colours for this slide"):
            pal = st.session_state.custom_palette or PALETTES[st.session_state.palette_name]
            sc1, sc2 = st.columns(2)
            with sc1:
                sp = st.color_picker("Primary", f"#{pal['primary']}", key=f"sp_{idx}")
            with sc2:
                sa = st.color_picker("Accent", f"#{pal['accent']}", key=f"sa_{idx}")
            st.session_state.setdefault("slide_color_overrides", {})[idx] = {
                "primary": sp.lstrip("#"), "accent": sa.lstrip("#")
            }

    with right:
        # Text editor
        st.markdown("#### Edit text")
        edits = st.session_state.slide_edits.get(idx, {})
        text_edits = edits.get("text_edits", {})

        for ti, text in enumerate(slide["texts"][:8]):
            new_text = st.text_input(f"Line {ti+1}", value=text, key=f"txt_{idx}_{ti}")
            if new_text != text:
                text_edits[text] = new_text

        if text_edits:
            st.session_state.slide_edits[idx] = {
                **edits,
                "text_edits": text_edits,
                "heading_size": st.session_state.get("cleanup_opts", {}).get("heading_size", 32),
                "body_size": st.session_state.get("cleanup_opts", {}).get("body_size", 16),
            }

        # AI design tips
        if ai and ai.get("design_tips"):
            st.markdown("#### Design tips")
            for tip in ai["design_tips"]:
                st.markdown(f"💡 {tip}")

        # AI typo fixes
        if ai and ai.get("typo_fixes"):
            st.markdown("#### Suggested fixes")
            for fix in ai["typo_fixes"]:
                col_a, col_b, col_c = st.columns([2,2,1])
                with col_a: st.markdown(f'~~{fix["original"]}~~')
                with col_b: st.markdown(f'**{fix["fixed"]}**')
                with col_c:
                    if st.button("Apply", key=f"fix_{idx}_{fix['original']}"):
                        text_edits[fix["original"]] = fix["fixed"]
                        st.session_state.slide_edits[idx] = {
                            **edits, "text_edits": text_edits
                        }
                        st.rerun()

    col_back, col_next = st.columns([1,5])
    with col_back:
        if st.button("← Back"): st.session_state.step = 4; st.rerun()
    with col_next:
        if st.button("Beautify my deck ✨", type="primary"): st.session_state.step = 6; st.rerun()

def step_beautify():
    st.markdown('<div class="main-header">Beautifying your deck…</div>', unsafe_allow_html=True)

    pal = st.session_state.custom_palette or PALETTES[st.session_state.palette_name]

    with st.spinner("Applying palette, fonts, edits and cleanup…"):
        result = apply_beautification(
            st.session_state.pptx_bytes,
            pal,
            st.session_state.font_pair,
            st.session_state.slide_templates,
            st.session_state.slide_edits,
        )
        st.session_state.beautified_pptx = result
        st.session_state.beautified = True

    st.success("🎨 Your deck has been Decorated by Deckorate AI!")

    # Summary of changes
    st.markdown("#### What changed")
    opts = st.session_state.get("cleanup_opts", {})
    heading_font, body_font = FONT_PAIRS[st.session_state.font_pair]
    changes = [
        f"Palette: **{st.session_state.palette_name}**",
        f"Heading font: **{heading_font}**",
        f"Body font: **{body_font}**",
        f"Heading size: **{opts.get('heading_size',32)}pt**",
        f"Body size: **{opts.get('body_size',16)}pt**",
        f"Slides with custom edits: **{len(st.session_state.slide_edits)}**",
        f"Slides with custom layouts: **{len(st.session_state.slide_templates)}**",
    ]
    for c in changes:
        st.markdown(f"• {c}")

    if st.button("Go to download →", type="primary"):
        st.session_state.step = 7
        st.rerun()

def step_download():
    st.markdown('<div class="main-header">🎨 Download your deck</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">Your Deckorate AI deck is ready. Fully editable in PowerPoint or Google Slides.</div>', unsafe_allow_html=True)

    if st.session_state.beautified_pptx:
        st.download_button(
            label="⬇️ Download your Deckorate'd .pptx",
            data=st.session_state.beautified_pptx,
            file_name="deckorate_ai_deck.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="primary",
            use_container_width=True,
        )

        st.markdown("---")
        st.markdown("#### Start over with a new deck?")
        if st.button("↺ Upload a new deck"):
            for key in list(st.session_state.keys()):
                del st.session_state[key]
            init_state()
            st.rerun()
    else:
        st.error("No beautified deck found. Please go back and run beautification.")
        if st.button("← Back to beautify"):
            st.session_state.step = 6
            st.rerun()

# ─── Main router ──────────────────────────────────────────────────────────────
render_sidebar()

step_map = {
    1: step_upload,
    2: step_palette,
    3: step_fonts,
    4: step_cleanup,
    5: step_preview,
    6: step_beautify,
    7: step_download,
}

step_map[st.session_state.step]()
